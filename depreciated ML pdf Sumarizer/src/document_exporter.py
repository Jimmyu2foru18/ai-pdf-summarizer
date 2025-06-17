import os
from pathlib import Path
from typing import Dict, List, Any
import io

# For Word document generation
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH

# For PowerPoint generation
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN

class DocumentExporter:
    """Class for exporting summaries and examples to various document formats."""
    
    def __init__(self):
        """Initialize the document exporter."""
        pass
    
    def export_to_word(self, document_structure: Dict, chapter_summaries: Dict, 
                      topic_summaries: Dict, topic_examples: Dict) -> io.BytesIO:
        """Export summaries and examples to a Word document.
        
        Args:
            document_structure: Structure of the document with chapters and topics
            chapter_summaries: Dictionary of chapter summaries
            topic_summaries: Dictionary of topic summaries
            topic_examples: Dictionary of topic examples
            
        Returns:
            BytesIO object containing the Word document
        """
        # Create a new Word document
        doc = Document()
        
        # Add title
        title = doc.add_heading('AI PDF Textbook Summarizer - Generated Summaries', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add chapter summaries
        doc.add_heading('Chapter Summaries', 1)
        for chapter_id, chapter_data in document_structure.items():
            doc.add_heading(f"Chapter {chapter_id}: {chapter_data['title']}", 2)
            doc.add_paragraph(chapter_summaries[chapter_id])
            doc.add_paragraph()
        
        # Add topic summaries
        doc.add_heading('Topic Summaries', 1)
        for chapter_id, chapter_data in document_structure.items():
            doc.add_heading(f"Chapter {chapter_id}: {chapter_data['title']}", 2)
            for topic_id, topic_data in chapter_data['topics'].items():
                doc.add_heading(f"{topic_data['title']}", 3)
                doc.add_paragraph(topic_summaries[topic_id])
                doc.add_paragraph()
        
        # Add examples
        doc.add_heading('Examples', 1)
        for chapter_id, chapter_data in document_structure.items():
            doc.add_heading(f"Chapter {chapter_id}: {chapter_data['title']}", 2)
            for topic_id, topic_data in chapter_data['topics'].items():
                doc.add_heading(f"Examples for: {topic_data['title']}", 3)
                examples = topic_examples[topic_id]
                for i, example in enumerate(examples, 1):
                    p = doc.add_paragraph()
                    p.add_run(f"Example {i}:").bold = True
                    doc.add_paragraph(example)
                doc.add_paragraph()
        
        # Save to a BytesIO object
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        return file_stream
    
    def export_to_powerpoint(self, document_structure: Dict, chapter_summaries: Dict, 
                           topic_summaries: Dict, topic_examples: Dict) -> io.BytesIO:
        """Export summaries and examples to a PowerPoint presentation.
        
        Args:
            document_structure: Structure of the document with chapters and topics
            chapter_summaries: Dictionary of chapter summaries
            topic_summaries: Dictionary of topic summaries
            topic_examples: Dictionary of topic examples
            
        Returns:
            BytesIO object containing the PowerPoint presentation
        """
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "AI PDF Textbook Summarizer"
        subtitle.text = "Generated Summaries and Examples"
        
        # Add section slide for chapter summaries
        section_slide_layout = prs.slide_layouts[2]  # Section header layout
        slide = prs.slides.add_slide(section_slide_layout)
        title = slide.shapes.title
        title.text = "Chapter Summaries"
        
        # Add chapter summary slides
        content_slide_layout = prs.slide_layouts[1]  # Content layout with title and content
        for chapter_id, chapter_data in document_structure.items():
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"Chapter {chapter_id}: {chapter_data['title']}"
            content.text = chapter_summaries[chapter_id]
        
        # Add section slide for topic summaries
        slide = prs.slides.add_slide(section_slide_layout)
        title = slide.shapes.title
        title.text = "Topic Summaries"
        
        # Add topic summary slides
        for chapter_id, chapter_data in document_structure.items():
            # Add chapter header slide
            slide = prs.slides.add_slide(section_slide_layout)
            title = slide.shapes.title
            title.text = f"Chapter {chapter_id}: {chapter_data['title']}"
            
            # Add topic slides for this chapter
            for topic_id, topic_data in chapter_data['topics'].items():
                slide = prs.slides.add_slide(content_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = topic_data['title']
                content.text = topic_summaries[topic_id]
        
        # Add section slide for examples
        slide = prs.slides.add_slide(section_slide_layout)
        title = slide.shapes.title
        title.text = "Examples"
        
        # Add example slides
        for chapter_id, chapter_data in document_structure.items():
            # Add chapter header slide for examples
            slide = prs.slides.add_slide(section_slide_layout)
            title = slide.shapes.title
            title.text = f"Chapter {chapter_id}: {chapter_data['title']} - Examples"
            
            # Add example slides for each topic
            for topic_id, topic_data in chapter_data['topics'].items():
                examples = topic_examples[topic_id]
                
                # Create a slide for this topic's examples
                slide = prs.slides.add_slide(content_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = f"Examples for: {topic_data['title']}"
                
                # Format examples text
                examples_text = ""
                for i, example in enumerate(examples, 1):
                    examples_text += f"Example {i}:\n{example}\n\n"
                
                content.text = examples_text
        
        # Save to a BytesIO object
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        return file_stream